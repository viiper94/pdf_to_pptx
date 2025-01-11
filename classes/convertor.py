import os
import time
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches


class Convertor:

    def __init__(self, index, file, thread, settings):
        self.thread = thread
        self.file = file
        self.index = index
        self.settings = settings

        self.slide_width = Inches(16)
        self.slide_height = Inches(self.get_height_multiplier(width=file.width, height=file.height))
        self.slide_aspect = self.slide_width / self.slide_height

        self.thread.file_process_start.emit(index)

        self.create_tmp_dir()

    def convert(self):
        try:
            # Saving start timestamp
            start = time.time()

            file_path = None

            if self.settings.output == 'pptx':
                prs = self.create_new_pptx_file()

                # Add slides
                for page_number in range(self.file.slides):
                    self.thread.file_process_progress.emit(self.index, page_number + 1, self.file.slides)
                    page = self.file.doc[page_number]
                    pil_image = self.page_to_pil(page)
                    slide = self.create_new_slide(prs)
                    self.insert_image_to_slide(slide, pil_image)

                file_path = self.save_pptx(prs)
            else:
                self.create_images_output_dir()

                # Save images
                for page_number in range(self.file.slides):
                    self.thread.file_process_progress.emit(self.index, page_number + 1, self.file.slides)

                    page = self.file.doc[page_number]
                    pil = self.page_to_pil(page)
                    file_path = f"{self.file.path_no_ext}/{page_number + 1}.png"
                    pil.save(file_path, 'PNG')

            # Saving end timestamp
            end = time.time()
            time_spent = end - start

            self.thread.file_process_end.emit(self.index, time_spent, file_path)

        except Exception as e:
            print(e)
            self.thread.file_process_failed.emit(self.index, "Виникла помилка =(")

        return True

    def create_new_pptx_file(self):
        prs = Presentation(self.settings.get_template_path())
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        return prs

    def get_slide_scale(self, page_width):
        # Calculate the scale factor for the desired resolution
        if self.settings.resolution:
            scale_x = self.settings.resolution / page_width
        else:
            scale_x = 1
        return scale_x

    def page_to_pil(self, page):
        bitmap = page.render(scale=self.get_slide_scale(page.get_width()))
        pil_image = bitmap.to_pil()
        return pil_image

    def pil_io_bytes(self, pil_image):
        image_bytes = BytesIO()
        pil_image.save(image_bytes, 'PNG')
        return image_bytes

    def page_to_file(self, page, file_path):
        bitmap = page.render(scale=self.get_slide_scale(page.get_width()))
        bitmap.save(file_path)
        return file_path

    def create_new_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        return slide

    def insert_image_to_slide(self, slide, pil_image):
        aspect_ratio = float(pil_image.width) / float(pil_image.height)
        image_bytes = self.pil_io_bytes(pil_image)
        # if image is wider than slide
        if aspect_ratio >= self.slide_aspect:
            height = int(self.slide_width / aspect_ratio)
            slide.shapes.add_picture(
                image_bytes,
                left=int(0),
                top=int((self.slide_height - height) / 2),
                height=height,
                width=self.slide_width
            )
        # if image is higher than slide [ ->| |<- ]
        else:
            width = int(self.slide_height * aspect_ratio)
            slide.shapes.add_picture(
                image_bytes,
                left=int((self.slide_width - width) / 2),
                top=int(0),
                height=self.slide_height,
                width=width
            )

    def save_pptx(self, prs):
        file_path = self.file.path_no_ext + '.pptx'
        prs.save(file_path)
        return file_path

    def create_images_output_dir(self):
        # creating dir with images
        if not os.path.exists(self.file.path_no_ext):
            os.mkdir(self.file.path_no_ext)
        return self.file.path_no_ext

    def create_tmp_dir(self):
        # creating dir for processed pdf slides
        if not os.path.exists(self.settings.get_tmp_folder_path()):
            os.mkdir(self.settings.get_tmp_folder_path())

    def get_height_multiplier(self, width, height):
        if self.settings.aspect == 'auto':
            aspect = width / height
            return 16 / aspect
        if self.settings.aspect == '16x9':
            return 9
        if self.settings.aspect == '4x3':
            return 12
