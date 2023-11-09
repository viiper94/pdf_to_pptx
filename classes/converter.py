import sys
import os
import re
import time
from io import BytesIO
from pdf2image import convert_from_path
from pdf2image import pdfinfo_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class Converter:

    # Change path prefix for .exe
    path_prefix = sys._MEIPASS
    # path_prefix = '.'

    def __init__(self, index, file, gui):
        self.file = file
        self.index = index
        self.gui = gui

        self.poppler_path = self.path_prefix + '/lib/poppler/bin'
        self.tmp_path = self.path_prefix + '/tmp'
        self.template_path = self.path_prefix + '/template/default.pptx'

        Image.MAX_IMAGE_PIXELS = 1000000000
        self.file_name_without_ext = self.get_file_name_without_extension()
        self.cpu_threads = self.get_threads()

        pdf_data = self.get_pdf_metadata(file)
        self.pages = pdf_data['Pages']
        self.size = pdf_data['File size']
        self.page_size = pdf_data['Page size']

        self.gui.update_gui_on_file_process(index, file_size=self.size, file_pages=self.pages)

        self.create_tmp_dir()
        self.convert()

    def convert(self):
        # Saving start timestamp
        start = time.time()

        # Convert PDF to images
        images = convert_from_path(
            self.file,
            dpi=300,
            fmt='jpeg',
            poppler_path=self.poppler_path,
            output_folder=self.tmp_path,
            thread_count=self.cpu_threads,
            size=(None, 1440)
        )

        self.create_pptx_from_images(images)

        # Saving end timestamp
        end = time.time()
        time_spent = end - start

        self.gui.update_gui_on_end(self.index, time_spent=time_spent)

        return True

    def get_pdf_metadata(self, pdf_file):
        return pdfinfo_from_path(
            pdf_file,
            poppler_path=self.poppler_path
        )

    def get_threads(self):
        if os.cpu_count() > 1:
            cpu_count = os.cpu_count() - 1
        else:
            cpu_count = os.cpu_count()
        return cpu_count

    def get_file_name_without_extension(self):
        return os.path.splitext(self.file)[0]

    def create_tmp_dir(self):
        # creating dir for processed pdf slides
        if not os.path.exists(self.tmp_path):
            os.mkdir(self.tmp_path)

    def create_pptx_from_images(self, images):
        # Create a new PowerPoint presentation
        prs = Presentation(self.template_path)
        page_height = self.get_height_multiplier()
        slide_width = Inches(16)
        slide_height = Inches(page_height)
        slide_aspect = slide_width / slide_height
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        # Add slides
        for i, image in enumerate(images):
            self.gui.update_gui_on_convertion(self.index, i + 1, self.pages)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            image_binary = BytesIO()
            image.save(image_binary, 'PNG')
            aspect_ratio = float(image.width) / float(image.height)
            if aspect_ratio >= slide_aspect:
                height = int(slide_width / aspect_ratio)
                slide.shapes.add_picture(
                    image_binary,
                    left=int(0),
                    top=int((slide_height - height) / 2),
                    height=height,
                    width=slide_width
                )

            else:
                width = int(Inches(page_height) * aspect_ratio)
                slide.shapes.add_picture(
                    image_binary,
                    left=int((slide_width - width) / 2),
                    top=int(0),
                    height=slide_height,
                    width=width
                )

        # Save the PowerPoint presentation
        prs.save(self.file_name_without_ext + '.pptx')
        return True

    def get_height_multiplier(self):
        matches = re.findall(r'(\d+\.?\d+)', self.page_size)
        width = float(matches[0])
        height = float(matches[1])
        aspect = width / height
        return 16 / aspect
