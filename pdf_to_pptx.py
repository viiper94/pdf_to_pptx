import sys
import os
import re
import time
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from pdf2image import pdfinfo_from_path
from io import BytesIO
from colorama import init, Fore
init()

def get_height_multiplier(pdf):
    matches = re.findall(r'(\d+\.?\d+)', pdf['Page size'])
    width = float(matches[0])
    height = float(matches[1])
    aspect = width/height
    return 16/aspect

def is_pdf_file(file_path):
    _, extension = os.path.splitext(file_path)
    return extension.lower() == ".pdf"

def pdf_to_pptx(pdf_file):

    # Change path prefix for .exe
    path_prefix = sys._MEIPASS
    # path_prefix = '.'

    # Gathering some data
    start = time.time()
    cpu_count = os.cpu_count()
    pdf = pdfinfo_from_path(pdf_file, poppler_path=path_prefix+'/poppler')
    pages = pdf['Pages']
    size = pdf['File size']

    print(Fore.RESET + "____________________________________")
    print(f"Конвертуємо файл {pdf_file}")
    print(Fore.WHITE + "Розмір файлу {:.2f} MB".format(int(size.split()[0]) / 10 ** 6))
    print(Fore.WHITE + f"Всього слайдів: {pages}")
    print(Fore.YELLOW + "Треба трохи зачекати, йде обробка файлу...")
    print(Fore.YELLOW + "Це може зайняти певний час, не закривайте це вікно...")

    # Get the file name without the extension
    file_name = os.path.splitext(pdf_file)[0]

    # creating dir for processed pdf slides
    if not os.path.exists(path_prefix+'/tmp'):
        os.mkdir(path_prefix+'/tmp')

    # Convert PDF to images
    images = convert_from_path(pdf_file, 300, poppler_path=path_prefix+'/poppler', output_folder=path_prefix+'/tmp', thread_count=cpu_count)

    num_slides = len(images)

    print(Fore.RESET + "____________________________________")

    # Create a new PowerPoint presentation
    prs = Presentation(path_prefix+"/template/default.pptx")
    height = get_height_multiplier(pdf)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(height)

    # Add slides
    for i, image in enumerate(images):
        print(Fore.RED + f"Конвертуємо слайд {i+1} з {num_slides}", end="\r")
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        image_binary = BytesIO()
        image.save(image_binary, 'PNG')
        aspect_ratio = float(image.width) / float(image.height)
        if aspect_ratio > (16/9):
            height = int(Inches(16) / aspect_ratio)
            slide.shapes.add_picture(image_binary, 0, int((Inches(height) - height) / 2), height=height, width=Inches(16))
        else:
            width = int(Inches(height) * aspect_ratio)
            slide.shapes.add_picture(image_binary, int((Inches(16) - width) / 2), 0, height=Inches(height), width=width)

    # Save the PowerPoint presentation
    prs.save(file_name + '.pptx')

    # Saving end timestamp
    end = time.time()

    sys.stdout.flush()
    print(Fore.GREEN + f"Конвертування файлу {pdf_file} завершено!")
    print(Fore.WHITE + f"Витрачено часу: {end - start:.2f}с")
    
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Щоб використовувати введіть в консолі:")
        print(Fore.YELLOW + "pdf_to_pptx.exe <pdf_file> [, <pdf_file>]")
        print(Fore.RESET + "Або перетягніть документ(и) на іконку програми")
        time.sleep(3)
    else:
        for argv in sys.argv:

            # Skipping first argument
            if argv == sys.argv[0]:
                continue

            # Check if file is .pdf
            if not is_pdf_file(argv):
                print(Fore.RED + f"Файл {argv} не PDF документ. Пропускаємо!")
            else:
                pdf_to_pptx(argv)

        time.sleep(3)
